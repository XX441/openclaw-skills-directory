#!/usr/bin/env python3
"""
AI PPT Generator — Assembly Script

Usage:
    python3 assembly.py \\
        --slides slide-*.png \\
        --brand igso \\
        --output "Deck_Name_v1.0.pptx" \\
        --video-slide 3 \\
        --youtube-url "https://www.youtube.com/embed/VIDEO_ID"

Or as a module:
    from assembly import build_deck
    build_deck(slides=[...], brand="igso", output="deck.pptx")
"""

import argparse
import json
import os
import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─── Brand Config ──────────────────────────────────────────────────

def load_brand_config(brand_name: str) -> dict:
    """Load brand config from brands/<name>/config.json"""
    skill_dir = Path(__file__).parent
    config_path = skill_dir / "brands" / brand_name / "config.json"
    
    if not config_path.exists():
        print(f"⚠️  Brand '{brand_name}' config not found. Using IGSO defaults.")
        config_path = skill_dir / "brands" / "igso" / "config.json"
        if not config_path.exists():
            raise FileNotFoundError(f"No brand config found. Tried: {config_path}")
    
    with open(config_path) as f:
        return json.load(f)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex string (#273F20) to python-pptx RGBColor."""
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─── Deck Builder ──────────────────────────────────────────────────

def build_deck(
    slides: list,
    brand: str = "igso",
    output: str = "deck.pptx",
    video_slide: int = None,
    youtube_url: str = None,
    video_label: str = None,
):
    """
    Build a PPTX from slide image files.
    
    Args:
        slides: List of image file paths (order determines slide order).
                Use None for a video placeholder slide.
        brand: Brand config name (maps to brands/<name>/config.json).
        output: Output PPTX file path.
        video_slide: 1-indexed slide number to place video on (if None, images only).
        youtube_url: YouTube embed URL for the video slide.
        video_label: Text label to display below the video.
    """
    config = load_brand_config(brand)
    
    primary = hex_to_rgb(config["primary"])
    accent = hex_to_rgb(config["accent"])
    text_light = hex_to_rgb(config["text_light"])
    bg_dark = hex_to_rgb(config["bg_dark"])
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    temp_dir = None
    
    for i, slide_path in enumerate(slides):
        slide_num = i + 1
        
        if slide_path is None or (
            video_slide and slide_num == video_slide and youtube_url
        ):
            # Video slide
            if youtube_url:
                slide = _add_video_slide(prs, youtube_url, video_label, config)
            else:
                slide = _add_placeholder_slide(prs, config, video_label)
        else:
            # Image slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            if os.path.exists(str(slide_path)):
                slide.shapes.add_picture(
                    str(slide_path),
                    Inches(0), Inches(0),
                    Inches(13.333), Inches(7.5)
                )
            else:
                print(f"⚠️  Slide {slide_num}: image not found: {slide_path}")
                # Add error text
                txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
                p = txBox.text_frame.paragraphs[0]
                p.text = f"[Missing image: {slide_path}]"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    
    # Save
    prs.save(output)
    size_kb = os.path.getsize(output) / 1024
    print(f"✅ Deck saved: {output}")
    print(f"   Slides: {len(prs.slides)}, Size: {size_kb:.0f} KB")
    
    # YouTube embed requires ZIP-level XML manipulation
    if youtube_url and video_slide:
        _embed_youtube_in_pptx(output, youtube_url, config)
    
    return output


def _add_video_slide(prs, youtube_url, label, config):
    """Add a slide with YouTube video placeholder (green bg + play icon)."""
    primary = hex_to_rgb(config["primary"])
    text_light = hex_to_rgb(config["text_light"])
    accent = hex_to_rgb(config["accent"])
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = primary
    
    # Darker play area rectangle
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(1.5),
        Inches(9.333), Inches(4.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x15)
    rect.line.fill.background()
    
    # Video label
    display_text = label or "▶  Play Video"
    txBox = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(8.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"▶  {display_text}"
    p.font.size = Pt(36)
    p.font.color.rgb = text_light
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Plays directly in PowerPoint when clicked"
    p2.font.size = Pt(16)
    p2.font.color.rgb = accent
    p2.alignment = PP_ALIGN.CENTER
    
    return slide


def _add_placeholder_slide(prs, config, label):
    """Add a plain placeholder (fallback if no video URL)."""
    return _add_video_slide(prs, None, label, config)


def _embed_youtube_in_pptx(pptx_path, youtube_url, config):
    """
    Replace a slide in the PPTX with YouTube video embed XML.
    This is needed because python-pptx doesn't support online video natively.
    """
    from lxml import etree
    
    temp_path = pptx_path + ".tmp"
    updated = False
    
    # YouTube embed XML from the original working example
    # Build the namespace URIs with proper curly braces
    ns1 = '{FF2B5EF4-FFF2-40B4-BE49-F238E27FC236}'
    ns2 = '{B00DDD20-D4F1-6E42-2336-D1F133E12682}'
    ns3 = '{BB962C8B-B14F-4D97-AF65-F5344CB8AC3E}'
    
    slide_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:pic>
        <p:nvPicPr>
          <p:cNvPr id="4" name="Online Media 3" title="{config.get('name', 'Video')}">
            <a:hlinkClick r:id="" action="ppaction://media"/>
            <a:extLst>
              <a:ext uri="{ns1}">
                <a16:creationId xmlns:a16="http://schemas.microsoft.com/office/drawing/2014/main"
                                id="{ns2}"/>
              </a:ext>
            </a:extLst>
          </p:cNvPr>
          <p:cNvPicPr>
            <a:picLocks noRot="1" noChangeAspect="1"/>
          </p:cNvPicPr>
          <p:nvPr>
            <a:videoFile r:link="rId1"/>
          </p:nvPr>
        </p:nvPicPr>
        <p:blipFill>
          <a:blip r:embed="rId2"/>
          <a:stretch>
            <a:fillRect/>
          </a:stretch>
        </p:blipFill>
        <p:spPr>
          <a:xfrm>
            <a:off x="1870841" y="1010882"/>
            <a:ext cx="8531200" cy="4798800"/>
          </a:xfrm>
          <a:prstGeom prst="rect">
            <a:avLst/>
          </a:prstGeom>
        </p:spPr>
      </p:pic>
    </p:spTree>
    <p:extLst>
      <p:ext uri="{ns3}">
        <p14:creationId xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
                        val="284430919"/>
      </p:ext>
    </p:extLst>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>
</p:sld>'''
    
    # Namespace URIs already have correct curly braces via f-string variables
    
    # Rels XML for the video slide
    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                Target="../media/image2.png"/>
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/video"
                Target="{youtube_url}" TargetMode="External"/>
</Relationships>'''
    
    # We need to find which slide is the video slide
    # The PPTX was built by python-pptx with slides sequentially numbered
    # We need to figure out which slideNN.xml to replace
    
    with zipfile.ZipFile(pptx_path, 'r') as z_in:
        with zipfile.ZipFile(temp_path, 'w', zipfile.ZIP_DEFLATED) as z_out:
            for item in z_in.infolist():
                data = z_in.read(item.filename)
                
                # Check if this is a slide file
                slide_match = re.match(r'ppt/slides/slide(\d+)\.xml', item.filename)
                if slide_match:
                    slide_num = int(slide_match.group(1))
                    # The video slide is the one that corresponds to video_slide
                    # In the current deck, we need to map: python-pptx creates slides in order
                    # but we added the video at build time. Since we already saved the pptx
                    # with a placeholder, we need to replace the correct slide.
                    
                    # Check if this slide contains our video placeholder text
                    if b"Play" in data[:500] or b"Sell" in data[:500]:
                        print(f"  Replacing slide {slide_num} XML with YouTube embed")
                        z_out.writestr(item, slide_xml.encode('utf-8'))
                        updated = True
                        continue
                
                # Replace rels for the video slide
                rels_match = re.match(r'ppt/slides/_rels/slide(\d+)\.xml\.rels', item.filename)
                if rels_match:
                    slide_num = int(rels_match.group(1))
                    # Check if corresponding slide was replaced
                    slide_file = f'ppt/slides/slide{slide_num}.xml'
                    try:
                        z_in.getinfo(slide_file)
                    except KeyError:
                        z_out.writestr(item, data)
                        continue
                    
                    slide_content = z_in.read(slide_file)
                    if b"Play" in slide_content[:500] or b"Sell" in slide_content[:500]:
                        print(f"  Replacing slide {slide_num} rels with YouTube reference")
                        z_out.writestr(item, rels_xml.encode('utf-8'))
                        continue
                
                z_out.writestr(item, data)
    
    os.replace(temp_path, pptx_path)
    
    if updated:
        print(f"✅ YouTube embed applied: {youtube_url}")
    else:
        print(f"⚠️  Could not find video placeholder slide to replace with YouTube embed.")
        print(f"   The video URL is noted. Manual YouTube embed may be needed.")
    
    final_size = os.path.getsize(pptx_path) / 1024
    print(f"   Final size: {final_size:.0f} KB")


# ─── CLI ───────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Build PPTX from slide images")
    parser.add_argument("--slides", nargs="+", required=True,
                        help="Slide image files in order")
    parser.add_argument("--brand", default="igso",
                        help="Brand config name (default: igso)")
    parser.add_argument("--output", default="deck.pptx",
                        help="Output PPTX file")
    parser.add_argument("--video-slide", type=int, default=None,
                        help="1-indexed slide number for video")
    parser.add_argument("--youtube-url", default=None,
                        help="YouTube embed URL for video slide")
    parser.add_argument("--video-label", default=None,
                        help="Label text for video slide")
    
    args = parser.parse_args()
    
    build_deck(
        slides=args.slides,
        brand=args.brand,
        output=args.output,
        video_slide=args.video_slide,
        youtube_url=args.youtube_url,
        video_label=args.video_label,
    )


if __name__ == "__main__":
    main()
